#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint转图片幻灯片工具 - 图形化界面版本
提供简单易用的图形界面来转换PPT文件
"""

import os
import sys
import threading
import tempfile
import shutil
from pathlib import Path
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import win32com.client
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class PPTConverterGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint转图片幻灯片工具")
        self.root.geometry("800x600")
        self.root.resizable(True, True)
        
        # 设置图标（如果有的话）
        try:
            self.root.iconbitmap(default="")
        except:
            pass
        
        # 变量
        self.input_file = tk.StringVar()
        self.output_file = tk.StringVar()
        self.image_format = tk.StringVar(value="PNG")
        self.temp_dir = tk.StringVar()
        self.use_custom_temp = tk.BooleanVar(value=False)
        self.is_converting = False
        
        self.setup_ui()
        
    def setup_ui(self):
        """设置用户界面"""
        # 主框架
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # 配置网格权重
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.rowconfigure(6, weight=1)
        
        # 标题
        title_label = ttk.Label(main_frame, text="PowerPoint转图片幻灯片工具", 
                               font=("", 16, "bold"))
        title_label.grid(row=0, column=0, columnspan=3, pady=(0, 20))
        
        # 输入文件选择
        ttk.Label(main_frame, text="输入PPT文件:").grid(row=1, column=0, sticky=tk.W, pady=5)
        ttk.Entry(main_frame, textvariable=self.input_file, width=50).grid(
            row=1, column=1, sticky=(tk.W, tk.E), padx=(10, 5), pady=5)
        ttk.Button(main_frame, text="浏览...", 
                  command=self.browse_input_file).grid(row=1, column=2, pady=5)
        
        # 输出文件选择
        ttk.Label(main_frame, text="输出PPT文件:").grid(row=2, column=0, sticky=tk.W, pady=5)
        ttk.Entry(main_frame, textvariable=self.output_file, width=50).grid(
            row=2, column=1, sticky=(tk.W, tk.E), padx=(10, 5), pady=5)
        ttk.Button(main_frame, text="浏览...", 
                  command=self.browse_output_file).grid(row=2, column=2, pady=5)
        
        # 选项框架
        options_frame = ttk.LabelFrame(main_frame, text="转换选项", padding="10")
        options_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=10)
        options_frame.columnconfigure(1, weight=1)
        
        # 图片格式选择
        ttk.Label(options_frame, text="图片格式:").grid(row=0, column=0, sticky=tk.W, pady=5)
        format_frame = ttk.Frame(options_frame)
        format_frame.grid(row=0, column=1, sticky=tk.W, padx=(10, 0), pady=5)
        ttk.Radiobutton(format_frame, text="PNG (推荐)", variable=self.image_format, 
                       value="PNG").pack(side=tk.LEFT, padx=(0, 20))
        ttk.Radiobutton(format_frame, text="JPG", variable=self.image_format, 
                       value="JPG").pack(side=tk.LEFT)
        
        # 临时目录选项
        ttk.Checkbutton(options_frame, text="使用自定义临时目录", 
                       variable=self.use_custom_temp,
                       command=self.toggle_temp_dir).grid(row=1, column=0, columnspan=2, 
                                                         sticky=tk.W, pady=5)
        
        self.temp_dir_frame = ttk.Frame(options_frame)
        self.temp_dir_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=5)
        self.temp_dir_frame.columnconfigure(1, weight=1)
        
        ttk.Label(self.temp_dir_frame, text="临时目录:").grid(row=0, column=0, sticky=tk.W)
        self.temp_dir_entry = ttk.Entry(self.temp_dir_frame, textvariable=self.temp_dir, 
                                       width=40, state='disabled')
        self.temp_dir_entry.grid(row=0, column=1, sticky=(tk.W, tk.E), padx=(10, 5))
        self.temp_dir_button = ttk.Button(self.temp_dir_frame, text="浏览...", 
                                         command=self.browse_temp_dir, state='disabled')
        self.temp_dir_button.grid(row=0, column=2)
        
        # 转换按钮
        button_frame = ttk.Frame(main_frame)
        button_frame.grid(row=4, column=0, columnspan=3, pady=20)
        
        self.convert_button = ttk.Button(button_frame, text="开始转换", 
                                        command=self.start_conversion,
                                        style="Accent.TButton")
        self.convert_button.pack(side=tk.LEFT, padx=(0, 10))
        
        self.clear_button = ttk.Button(button_frame, text="清空", 
                                      command=self.clear_fields)
        self.clear_button.pack(side=tk.LEFT, padx=(0, 10))
        
        self.exit_button = ttk.Button(button_frame, text="退出", 
                                     command=self.root.quit)
        self.exit_button.pack(side=tk.LEFT)
        
        # 进度条
        self.progress_frame = ttk.Frame(main_frame)
        self.progress_frame.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=10)
        self.progress_frame.columnconfigure(0, weight=1)
        
        self.progress_var = tk.StringVar(value="就绪")
        ttk.Label(self.progress_frame, textvariable=self.progress_var).grid(
            row=0, column=0, sticky=tk.W)
        
        self.progress_bar = ttk.Progressbar(self.progress_frame, mode='indeterminate')
        self.progress_bar.grid(row=1, column=0, sticky=(tk.W, tk.E), pady=(5, 0))
        
        # 日志输出区域
        log_frame = ttk.LabelFrame(main_frame, text="转换日志", padding="5")
        log_frame.grid(row=6, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=10)
        log_frame.columnconfigure(0, weight=1)
        log_frame.rowconfigure(0, weight=1)
        
        self.log_text = scrolledtext.ScrolledText(log_frame, height=12, width=80)
        self.log_text.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # 添加欢迎信息
        self.log_message("欢迎使用PowerPoint转图片幻灯片工具！")
        self.log_message("请选择要转换的PPT文件。")
        
    def toggle_temp_dir(self):
        """切换临时目录选项"""
        if self.use_custom_temp.get():
            self.temp_dir_entry.config(state='normal')
            self.temp_dir_button.config(state='normal')
        else:
            self.temp_dir_entry.config(state='disabled')
            self.temp_dir_button.config(state='disabled')
            self.temp_dir.set("")
    
    def browse_input_file(self):
        """浏览输入文件"""
        filename = filedialog.askopenfilename(
            title="选择PowerPoint文件",
            filetypes=[
                ("PowerPoint文件", "*.pptx *.ppt"),
                ("所有文件", "*.*")
            ]
        )
        if filename:
            self.input_file.set(filename)
            # 自动设置输出文件名
            if not self.output_file.get():
                input_dir = os.path.dirname(filename)
                input_name = os.path.splitext(os.path.basename(filename))[0]
                output_path = os.path.join(input_dir, f"{input_name}_images.pptx")
                self.output_file.set(output_path)
    
    def browse_output_file(self):
        """浏览输出文件"""
        filename = filedialog.asksaveasfilename(
            title="保存图片PPT文件",
            defaultextension=".pptx",
            filetypes=[
                ("PowerPoint文件", "*.pptx"),
                ("所有文件", "*.*")
            ]
        )
        if filename:
            self.output_file.set(filename)
    
    def browse_temp_dir(self):
        """浏览临时目录"""
        dirname = filedialog.askdirectory(title="选择临时目录")
        if dirname:
            self.temp_dir.set(dirname)
    
    def clear_fields(self):
        """清空所有字段"""
        self.input_file.set("")
        self.output_file.set("")
        self.temp_dir.set("")
        self.use_custom_temp.set(False)
        self.image_format.set("PNG")
        self.toggle_temp_dir()
        self.log_text.delete(1.0, tk.END)
        self.log_message("字段已清空，请重新选择文件。")
    
    def log_message(self, message):
        """添加日志消息"""
        self.log_text.insert(tk.END, f"{message}\n")
        self.log_text.see(tk.END)
        self.root.update_idletasks()
    
    def validate_inputs(self):
        """验证输入"""
        if not self.input_file.get():
            messagebox.showerror("错误", "请选择输入PPT文件")
            return False
        
        if not os.path.exists(self.input_file.get()):
            messagebox.showerror("错误", "输入文件不存在")
            return False
        
        if not self.output_file.get():
            messagebox.showerror("错误", "请设置输出文件路径")
            return False
        
        if self.use_custom_temp.get() and not self.temp_dir.get():
            messagebox.showerror("错误", "请选择临时目录")
            return False
        
        return True
    
    def start_conversion(self):
        """开始转换"""
        if not self.validate_inputs():
            return
        
        if self.is_converting:
            messagebox.showwarning("警告", "转换正在进行中，请等待完成")
            return
        
        # 在新线程中运行转换
        self.is_converting = True
        self.convert_button.config(state='disabled')
        self.progress_bar.start()
        self.progress_var.set("转换中...")
        
        thread = threading.Thread(target=self.convert_ppt)
        thread.daemon = True
        thread.start()
    
    def convert_ppt(self):
        """执行PPT转换"""
        try:
            self.log_message("=" * 50)
            self.log_message("开始转换PPT文件...")
            
            # 获取参数
            input_ppt = self.input_file.get()
            output_ppt = self.output_file.get()
            image_format = self.image_format.get()
            temp_dir = self.temp_dir.get() if self.use_custom_temp.get() else None
            
            self.log_message(f"输入文件: {input_ppt}")
            self.log_message(f"输出文件: {output_ppt}")
            self.log_message(f"图片格式: {image_format}")
            
            # 创建转换器
            converter = PPTToImageSlides(self.log_message)
            
            # 执行转换
            converter.convert_ppt_to_image_slides(
                input_ppt, output_ppt, temp_dir, image_format
            )
            
            # 转换完成
            self.root.after(0, self.conversion_complete, True, "转换成功完成！")
            
        except Exception as e:
            error_msg = f"转换失败: {str(e)}"
            self.root.after(0, self.conversion_complete, False, error_msg)
    
    def conversion_complete(self, success, message):
        """转换完成回调"""
        self.is_converting = False
        self.convert_button.config(state='normal')
        self.progress_bar.stop()
        
        if success:
            self.progress_var.set("转换完成")
            self.log_message(message)
            messagebox.showinfo("成功", message)
        else:
            self.progress_var.set("转换失败")
            self.log_message(message)
            messagebox.showerror("错误", message)


class PPTToImageSlides:
    def __init__(self, log_callback=None):
        self.powerpoint = None
        self.log_callback = log_callback or print
        
    def log(self, message):
        """记录日志"""
        self.log_callback(message)
        
    def __enter__(self):
        try:
            # 启动PowerPoint应用程序
            self.powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            self.log("PowerPoint COM接口初始化成功")
            return self
        except Exception as e:
            self.log(f"初始化PowerPoint COM接口失败: {e}")
            raise
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.powerpoint:
            try:
                self.powerpoint.Quit()
                self.log("PowerPoint COM接口已关闭")
            except:
                pass
    
    def export_slides_to_images(self, ppt_path, output_dir, image_format="PNG"):
        """导出幻灯片为图片"""
        try:
            # 确保输出目录存在
            os.makedirs(output_dir, exist_ok=True)
            
            # 打开PPT文件
            presentation = self.powerpoint.Presentations.Open(str(Path(ppt_path).absolute()))
            
            # 获取幻灯片数量和尺寸信息
            slide_count = presentation.Slides.Count
            self.log(f"找到 {slide_count} 张幻灯片")
            
            # 获取幻灯片的页面设置信息
            page_setup = presentation.PageSetup
            slide_width_points = page_setup.SlideWidth  # 以点为单位
            slide_height_points = page_setup.SlideHeight  # 以点为单位
            
            self.log(f"幻灯片尺寸: {slide_width_points} x {slide_height_points} 点")
            
            # 计算幻灯片的英寸尺寸（72点 = 1英寸）
            slide_width_inches = slide_width_points / 72.0
            slide_height_inches = slide_height_points / 72.0
            
            # 设置导出格式
            if image_format.upper() == "PNG":
                file_extension = ".png"
            else:
                file_extension = ".jpg"
            
            exported_files = []
            slide_info = None
            
            # 导出每一张幻灯片
            for i in range(1, slide_count + 1):
                slide = presentation.Slides(i)
                output_file = os.path.join(output_dir, f"slide_{i:03d}{file_extension}")
                
                try:
                    # 导出幻灯片为图片
                    slide.Export(output_file, image_format.upper())
                    exported_files.append(output_file)
                    self.log(f"导出幻灯片 {i}/{slide_count}: {os.path.basename(output_file)}")
                    
                    # 获取第一张图片的信息来计算实际的尺寸比例
                    if i == 1 and os.path.exists(output_file):
                        with Image.open(output_file) as img:
                            img_width_pixels, img_height_pixels = img.size
                            
                            # 尝试获取图片的DPI信息
                            dpi_info = img.info.get('dpi', None)
                            if dpi_info:
                                dpi_x, dpi_y = dpi_info
                                avg_dpi = (dpi_x + dpi_y) / 2
                                self.log(f"从图片元数据读取DPI: X={dpi_x:.1f}, Y={dpi_y:.1f}, 平均={avg_dpi:.1f}")
                            else:
                                # 如果图片没有DPI信息，计算基于像素和点尺寸的比例
                                scale_factor_x = img_width_pixels / slide_width_points
                                scale_factor_y = img_height_pixels / slide_height_points
                                avg_scale_factor = (scale_factor_x + scale_factor_y) / 2
                                
                                # 估算等效DPI
                                estimated_dpi = avg_scale_factor * 72
                                
                                self.log(f"图片尺寸: {img_width_pixels}x{img_height_pixels} 像素")
                                self.log(f"幻灯片尺寸: {slide_width_points}x{slide_height_points} 点")
                                self.log(f"缩放因子: X={scale_factor_x:.2f}, Y={scale_factor_y:.2f}")
                                self.log(f"估算的等效DPI: {estimated_dpi:.1f} (基于像素/点数比例)")
                                
                                dpi_x = dpi_y = avg_dpi = estimated_dpi
                            
                            slide_info = {
                                'slide_width_inches': slide_width_inches,
                                'slide_height_inches': slide_height_inches,
                                'slide_width_points': slide_width_points,
                                'slide_height_points': slide_height_points,
                                'img_width_pixels': img_width_pixels,
                                'img_height_pixels': img_height_pixels,
                                'dpi_x': dpi_x,
                                'dpi_y': dpi_y,
                                'avg_dpi': avg_dpi
                            }
                    
                except Exception as e:
                    self.log(f"导出幻灯片 {i} 失败: {e}")
                    continue
            
            # 关闭PPT文件
            presentation.Close()
            
            return exported_files, slide_info
            
        except Exception as e:
            self.log(f"导出幻灯片时发生错误: {e}")
            return [], None
    
    def create_image_ppt(self, image_files, output_ppt_path, slide_info=None):
        """创建图片PPT"""
        try:
            # 创建新的PPT演示文稿
            prs = Presentation()
            
            # 如果有slide_info，使用原始PPT的页面尺寸设置新PPT
            if slide_info:
                # 将点转换为EMU单位（English Metric Units）
                # 1 point = 12700 EMU
                slide_width_emu = int(slide_info['slide_width_points'] * 12700)
                slide_height_emu = int(slide_info['slide_height_points'] * 12700)
                
                # 设置幻灯片尺寸
                prs.slide_width = slide_width_emu
                prs.slide_height = slide_height_emu
                
                self.log(f"设置幻灯片尺寸为原始PPT尺寸: {slide_info['slide_width_points']} x {slide_info['slide_height_points']} 点")
                self.log(f"对应英寸尺寸: {slide_info['slide_width_inches']:.2f}\" x {slide_info['slide_height_inches']:.2f}\"")
                self.log(f"检测到的DPI: {slide_info['avg_dpi']:.1f}")
            else:
                self.log("警告: 未获取到原始PPT尺寸信息，使用默认尺寸")
            
            # 获取设置后的幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide_width_inches = slide_width / 914400  # EMU to inches (914400 EMU = 1 inch)
            slide_height_inches = slide_height / 914400
            
            self.log(f"最终幻灯片尺寸: {slide_width_inches:.2f}\" x {slide_height_inches:.2f}\"")
            
            for i, image_file in enumerate(image_files, 1):
                if not os.path.exists(image_file):
                    self.log(f"图片文件不存在，跳过: {image_file}")
                    continue
                
                try:
                    # 添加空白幻灯片
                    slide_layout = prs.slide_layouts[6]  # 空白布局
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 获取图片尺寸
                    with Image.open(image_file) as img:
                        img_width_pixels, img_height_pixels = img.size
                    
                    # 计算图片应该显示的尺寸
                    if slide_info and 'avg_dpi' in slide_info:
                        # 使用检测到的DPI信息
                        img_width_inches = img_width_pixels / slide_info['avg_dpi']
                        img_height_inches = img_height_pixels / slide_info['avg_dpi']
                        if i == 1:
                            self.log(f"使用检测到的DPI {slide_info['avg_dpi']:.1f} 进行尺寸计算")
                            self.log(f"图片 {img_width_pixels}x{img_height_pixels} 像素 → {img_width_inches:.2f}x{img_height_inches:.2f} 英寸")
                    else:
                        # 如果没有DPI信息，直接按比例缩放使图片填满幻灯片
                        img_width_inches = slide_width_inches
                        img_height_inches = slide_height_inches
                        if i == 1:
                            self.log("无DPI信息，使用直接比例缩放")
                    
                    # 计算缩放比例，确保图片填满整个幻灯片
                    scale_x = slide_width_inches / img_width_inches
                    scale_y = slide_height_inches / img_height_inches
                    scale = max(scale_x, scale_y)  # 使用较大的缩放比例确保填满
                    
                    # 计算最终尺寸
                    final_width = Inches(img_width_inches * scale)
                    final_height = Inches(img_height_inches * scale)
                    
                    # 计算居中位置
                    left = (slide_width - final_width) / 2
                    top = (slide_height - final_height) / 2
                    
                    # 插入图片
                    slide.shapes.add_picture(
                        image_file, 
                        left, 
                        top, 
                        final_width, 
                        final_height
                    )
                    
                    self.log(f"添加图片到幻灯片 {i}: {os.path.basename(image_file)}")
                    
                except Exception as e:
                    self.log(f"处理图片 {image_file} 时发生错误: {e}")
                    continue
            
            # 确保输出目录存在
            output_dir = os.path.dirname(output_ppt_path)
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
            
            # 保存PPT文件
            prs.save(output_ppt_path)
            self.log(f"成功创建图片PPT文件: {output_ppt_path}")
            
        except Exception as e:
            self.log(f"创建图片PPT时发生错误: {e}")
            raise
    
    def convert_ppt_to_image_slides(self, input_ppt, output_ppt, temp_dir=None, image_format="PNG"):
        """主转换函数"""
        # 验证输入文件
        if not os.path.exists(input_ppt):
            raise FileNotFoundError(f"输入文件不存在: {input_ppt}")
        
        # 创建临时目录
        if temp_dir is None:
            temp_dir = tempfile.mkdtemp(prefix="ppt_to_image_")
            cleanup_temp = True
        else:
            os.makedirs(temp_dir, exist_ok=True)
            cleanup_temp = False
        
        try:
            self.log(f"临时目录: {temp_dir}")
            
            with self:
                # 导出幻灯片为图片
                image_files, slide_info = self.export_slides_to_images(input_ppt, temp_dir, image_format)
                
                if not image_files:
                    raise RuntimeError("没有成功导出任何图片")
                
                self.log(f"成功导出 {len(image_files)} 张图片")
                
                # 创建图片PPT
                self.create_image_ppt(image_files, output_ppt, slide_info)
                
                self.log("转换完成！")
            
        finally:
            # 清理临时文件
            if cleanup_temp and os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                    self.log(f"清理临时目录: {temp_dir}")
                except Exception as e:
                    self.log(f"清理临时目录失败: {e}")


def main():
    # 创建主窗口
    root = tk.Tk()
    
    # 设置主题
    try:
        style = ttk.Style()
        # 尝试使用现代主题
        available_themes = style.theme_names()
        if 'vista' in available_themes:
            style.theme_use('vista')
        elif 'clam' in available_themes:
            style.theme_use('clam')
    except:
        pass
    
    # 创建应用程序
    app = PPTConverterGUI(root)
    
    # 运行主循环
    root.mainloop()


if __name__ == "__main__":
    main()
