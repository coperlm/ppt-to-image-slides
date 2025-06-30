#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint转图片幻灯片脚本
功能：将PPT文件的每一页导出为图片，然后创建新的纯图片PPT文件
"""

import os
import sys
import argparse
import tempfile
import shutil
from pathlib import Path
import win32com.client
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class PPTToImageSlides:
    def __init__(self):
        self.powerpoint = None
        
    def __enter__(self):
        try:
            # 启动PowerPoint应用程序
            self.powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # 注意：某些PowerPoint版本不允许隐藏窗口，所以我们不设置Visible属性
            print("PowerPoint COM接口初始化成功")
            return self
        except Exception as e:
            print(f"初始化PowerPoint COM接口失败: {e}")
            raise
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.powerpoint:
            try:
                self.powerpoint.Quit()
                print("PowerPoint COM接口已关闭")
            except:
                pass
    
    def export_slides_to_images(self, ppt_path, output_dir, image_format="PNG"):
        """
        将PPT文件的每一页导出为图片
        
        Args:
            ppt_path (str): PPT文件路径
            output_dir (str): 图片输出目录
            image_format (str): 图片格式，PNG或JPG
        
        Returns:
            tuple: (导出的图片文件路径列表, 幻灯片尺寸信息)
        """
        try:
            # 确保输出目录存在
            os.makedirs(output_dir, exist_ok=True)
            
            # 打开PPT文件
            presentation = self.powerpoint.Presentations.Open(str(Path(ppt_path).absolute()))
            
            # 获取幻灯片数量和尺寸信息
            slide_count = presentation.Slides.Count
            print(f"找到 {slide_count} 张幻灯片")
            
            # 获取幻灯片的页面设置信息
            page_setup = presentation.PageSetup
            slide_width_points = page_setup.SlideWidth  # 以点为单位
            slide_height_points = page_setup.SlideHeight  # 以点为单位
            
            print(f"幻灯片尺寸: {slide_width_points} x {slide_height_points} 点")
            
            # 计算幻灯片的英寸尺寸（72点 = 1英寸）
            slide_width_inches = slide_width_points / 72.0
            slide_height_inches = slide_height_points / 72.0
            
            # 设置导出格式
            if image_format.upper() == "PNG":
                export_format = 18  # ppSaveAsPNG
                file_extension = ".png"
            else:
                export_format = 17  # ppSaveAsJPEG
                file_extension = ".jpg"
            
            exported_files = []
            slide_info = None
            
            # 导出每一张幻灯片
            for i in range(1, slide_count + 1):
                slide = presentation.Slides(i)
                output_file = os.path.join(output_dir, f"slide_{i:03d}{file_extension}")
                
                try:
                    # 导出幻灯片为图片
                    slide.Export(output_file, image_format.upper())
                    exported_files.append(output_file)
                    print(f"导出幻灯片 {i}/{slide_count}: {output_file}")
                    
                    # 获取第一张图片的信息来计算实际的尺寸比例
                    if i == 1 and os.path.exists(output_file):
                        with Image.open(output_file) as img:
                            img_width_pixels, img_height_pixels = img.size
                            
                            # 尝试获取图片的DPI信息
                            dpi_info = img.info.get('dpi', None)
                            if dpi_info:
                                dpi_x, dpi_y = dpi_info
                                avg_dpi = (dpi_x + dpi_y) / 2
                                print(f"图片包含DPI信息: X={dpi_x:.1f}, Y={dpi_y:.1f}, 平均={avg_dpi:.1f}")
                            else:
                                # 如果图片没有DPI信息，计算基于像素和点尺寸的比例
                                # 但不称其为DPI，而是作为缩放因子
                                scale_factor_x = img_width_pixels / slide_width_points
                                scale_factor_y = img_height_pixels / slide_height_points
                                avg_scale_factor = (scale_factor_x + scale_factor_y) / 2
                                
                                # 为了兼容性，我们假设这是基于某个DPI的结果
                                # 通常PowerPoint导出时会使用较高的分辨率
                                estimated_dpi = avg_scale_factor * 72  # 72是点到英寸的转换因子
                                
                                print(f"图片尺寸: {img_width_pixels}x{img_height_pixels} 像素")
                                print(f"幻灯片尺寸: {slide_width_points}x{slide_height_points} 点")
                                print(f"缩放因子: X={scale_factor_x:.2f}, Y={scale_factor_y:.2f}")
                                print(f"计算的等效DPI: {estimated_dpi:.1f}")
                                
                                dpi_x = dpi_y = avg_dpi = estimated_dpi
                            
                            slide_info = {
                                'slide_width_inches': slide_width_inches,
                                'slide_height_inches': slide_height_inches,
                                'slide_width_points': slide_width_points,
                                'slide_height_points': slide_height_points,
                                'img_width_pixels': img_width_pixels,
                                'img_height_pixels': img_height_pixels,
                                'dpi_x': dpi_x,
                                'dpi_y': dpi_y,
                                'avg_dpi': avg_dpi
                            }
                    
                except Exception as e:
                    print(f"导出幻灯片 {i} 失败: {e}")
                    continue
            
            # 关闭PPT文件
            presentation.Close()
            
            return exported_files, slide_info
            
        except Exception as e:
            print(f"导出幻灯片时发生错误: {e}")
            return [], None
    
    def create_image_ppt(self, image_files, output_ppt_path, slide_info=None):
        """
        创建由图片组成的新PPT文件
        
        Args:
            image_files (list): 图片文件路径列表
            output_ppt_path (str): 输出PPT文件路径
            slide_info (dict): 幻灯片尺寸和DPI信息
        """
        try:
            # 创建新的PPT演示文稿
            prs = Presentation()
            
            # 如果有slide_info，使用原始PPT的页面尺寸设置新PPT
            if slide_info:
                # 将点转换为EMU单位（English Metric Units）
                # 1 point = 12700 EMU
                slide_width_emu = int(slide_info['slide_width_points'] * 12700)
                slide_height_emu = int(slide_info['slide_height_points'] * 12700)
                
                # 设置幻灯片尺寸
                prs.slide_width = slide_width_emu
                prs.slide_height = slide_height_emu
                
                print(f"设置幻灯片尺寸为原始PPT尺寸: {slide_info['slide_width_points']} x {slide_info['slide_height_points']} 点")
                print(f"对应英寸尺寸: {slide_info['slide_width_inches']:.2f}\" x {slide_info['slide_height_inches']:.2f}\"")
                print(f"检测到的DPI: {slide_info['avg_dpi']:.1f}")
            else:
                print("警告: 未获取到原始PPT尺寸信息，使用默认尺寸")
            
            # 获取设置后的幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide_width_inches = slide_width / 914400  # EMU to inches (914400 EMU = 1 inch)
            slide_height_inches = slide_height / 914400
            
            print(f"最终幻灯片尺寸: {slide_width_inches:.2f}\" x {slide_height_inches:.2f}\"")
            
            for i, image_file in enumerate(image_files, 1):
                if not os.path.exists(image_file):
                    print(f"图片文件不存在，跳过: {image_file}")
                    continue
                
                try:
                    # 添加空白幻灯片
                    slide_layout = prs.slide_layouts[6]  # 空白布局
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 获取图片尺寸
                    with Image.open(image_file) as img:
                        img_width_pixels, img_height_pixels = img.size
                    
                    # 计算图片应该显示的尺寸
                    if slide_info and 'avg_dpi' in slide_info:
                        # 使用检测到的DPI信息
                        img_width_inches = img_width_pixels / slide_info['avg_dpi']
                        img_height_inches = img_height_pixels / slide_info['avg_dpi']
                        if i == 1:
                            print(f"使用检测到的DPI {slide_info['avg_dpi']:.1f} 进行尺寸计算")
                            print(f"图片 {img_width_pixels}x{img_height_pixels} 像素 → {img_width_inches:.2f}x{img_height_inches:.2f} 英寸")
                    else:
                        # 如果没有DPI信息，直接按比例缩放使图片填满幻灯片
                        # 这是最安全的方法
                        img_width_inches = slide_width_inches
                        img_height_inches = slide_height_inches
                        if i == 1:
                            print("无DPI信息，使用直接比例缩放")
                    
                    # 计算缩放比例，确保图片填满整个幻灯片
                    scale_x = slide_width_inches / img_width_inches
                    scale_y = slide_height_inches / img_height_inches
                    scale = max(scale_x, scale_y)  # 使用较大的缩放比例确保填满
                    
                    # 计算最终尺寸
                    final_width = Inches(img_width_inches * scale)
                    final_height = Inches(img_height_inches * scale)
                    
                    # 计算居中位置
                    left = (slide_width - final_width) / 2
                    top = (slide_height - final_height) / 2
                    
                    # 插入图片
                    slide.shapes.add_picture(
                        image_file, 
                        left, 
                        top, 
                        final_width, 
                        final_height
                    )
                    
                    print(f"添加图片到幻灯片 {i}: {os.path.basename(image_file)}")
                    
                except Exception as e:
                    print(f"处理图片 {image_file} 时发生错误: {e}")
                    continue
            
            # 确保输出目录存在
            output_dir = os.path.dirname(output_ppt_path)
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
            
            # 保存PPT文件
            prs.save(output_ppt_path)
            print(f"成功创建图片PPT文件: {output_ppt_path}")
            
        except Exception as e:
            print(f"创建图片PPT时发生错误: {e}")
            raise
    
    def convert_ppt_to_image_slides(self, input_ppt, output_ppt, temp_dir=None, image_format="PNG"):
        """
        主转换函数：将PPT转换为纯图片PPT
        
        Args:
            input_ppt (str): 输入PPT文件路径
            output_ppt (str): 输出PPT文件路径
            temp_dir (str): 临时目录，如果为None则自动创建
            image_format (str): 图片格式
        """
        # 验证输入文件
        if not os.path.exists(input_ppt):
            raise FileNotFoundError(f"输入文件不存在: {input_ppt}")
        
        # 创建临时目录
        if temp_dir is None:
            temp_dir = tempfile.mkdtemp(prefix="ppt_to_image_")
            cleanup_temp = True
        else:
            os.makedirs(temp_dir, exist_ok=True)
            cleanup_temp = False
        
        try:
            print(f"开始转换: {input_ppt}")
            print(f"临时目录: {temp_dir}")
            
            # 导出幻灯片为图片
            image_files, slide_info = self.export_slides_to_images(input_ppt, temp_dir, image_format)
            
            if not image_files:
                raise RuntimeError("没有成功导出任何图片")
            
            print(f"成功导出 {len(image_files)} 张图片")
            
            # 创建图片PPT
            self.create_image_ppt(image_files, output_ppt, slide_info)
            
            print("转换完成！")
            
        finally:
            # 清理临时文件
            if cleanup_temp and os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                    print(f"清理临时目录: {temp_dir}")
                except Exception as e:
                    print(f"清理临时目录失败: {e}")


def main():
    parser = argparse.ArgumentParser(description="将PowerPoint文件转换为纯图片幻灯片")
    parser.add_argument("input", help="输入PPT文件路径")
    parser.add_argument("-o", "--output", help="输出PPT文件路径（默认为输入文件名_images.pptx）")
    parser.add_argument("-f", "--format", choices=["PNG", "JPG"], default="PNG", 
                       help="图片格式（默认PNG）")
    parser.add_argument("-t", "--temp-dir", help="临时目录（可选）")
    
    args = parser.parse_args()
    
    # 处理输入路径
    input_ppt = os.path.abspath(args.input)
    
    # 处理输出路径
    if args.output:
        output_ppt = os.path.abspath(args.output)
    else:
        input_dir = os.path.dirname(input_ppt)
        input_name = os.path.splitext(os.path.basename(input_ppt))[0]
        output_ppt = os.path.join(input_dir, f"{input_name}_images.pptx")
    
    try:
        with PPTToImageSlides() as converter:
            converter.convert_ppt_to_image_slides(
                input_ppt, 
                output_ppt, 
                args.temp_dir, 
                args.format
            )
    except Exception as e:
        print(f"转换失败: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
